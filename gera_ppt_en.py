from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL_ESCURO  = RGBColor(0x1A, 0x3A, 0x6B)
AZUL_MEDIO   = RGBColor(0x2B, 0x57, 0x9A)
AZUL_CLARO   = RGBColor(0xD6, 0xE4, 0xF7)
LARANJA      = RGBColor(0xE8, 0x7B, 0x00)
CINZA_ESCURO = RGBColor(0x2D, 0x2D, 0x2D)
CINZA_MEDIO  = RGBColor(0x55, 0x5F, 0x6E)
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_BG     = RGBColor(0xF4, 0xF6, 0xFB)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]

import os
LOGO = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                    "NXProject.Community", "Assets", "branding", "logo-nexus-xdata-transparent.png")

def add_rect(slide, l, t, w, h, fill):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape

def add_textbox(slide, l, t, w, h, text, size, bold=False, color=CINZA_ESCURO,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, size, bold=False, color=CINZA_ESCURO,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def header(slide, titulo):
    add_rect(slide, 0, 0, W, H, CINZA_BG)
    add_rect(slide, 0, 0, W, Inches(1.3), AZUL_ESCURO)
    add_rect(slide, 0, Inches(1.3), Inches(0.07), H - Inches(1.3), LARANJA)
    add_textbox(slide, Inches(0.5), Inches(0.22), Inches(12), Inches(0.8),
                titulo, 28, bold=True, color=BRANCO)

# ── Slide 1 — Cover ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, AZUL_ESCURO)
add_rect(slide, 0, Inches(4.6), W, Inches(0.08), LARANJA)
add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(1.4),
            "NXProject", 72, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.0), Inches(11.3), Inches(0.7),
            "Intelligent planning integrated with Azure DevOps",
            24, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
            "From estimates to schedule — without leaving the technical backlog behind.",
            16, italic=True, color=RGBColor(0xB0, 0xC8, 0xF0), align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.4),
            "Nexus XData Tecnologia  •  nexusxdata.com.br",
            11, color=RGBColor(0x7A, 0x9A, 0xC8), align=PP_ALIGN.CENTER)

# ── Slide 2 — The problem ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "The problem NXProject solves")

dores = [
    ("📋  Rich backlog, missing schedule",
     "The technical team works in Azure DevOps, but management cannot\nsee dates, dependencies and allocation in one consolidated view."),
    ("⚠️  Rework and parallel spreadsheets",
     "Managers maintain Excel sheets manually, out of sync with what the\nteam is actually doing in the backlog."),
    ("🔗  Lack of traceability between planning and execution",
     "There is no direct link between the manager's estimate and the\ndeveloper's work item — every change requires double updates."),
    ("📅  Hard to replan quickly",
     "When a sprint slips, recalculating all dependencies and dates\nmanually is slow and error-prone."),
]
for i, (titulo, corpo) in enumerate(dores):
    col = i % 2; row = i // 2
    l = Inches(0.5 + col * 6.4); t = Inches(1.7 + row * 2.5)
    w = Inches(6.0); h = Inches(2.2)
    add_rect(slide, l, t, w, h, BRANCO)
    tb = slide.shapes.add_textbox(l+Inches(0.18), t+Inches(0.15), w-Inches(0.3), h-Inches(0.3))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = titulo
    run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = AZUL_ESCURO
    add_para(tf, corpo, 11, color=CINZA_MEDIO, space_before=6)

# ── Slide 3 — What it is ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "What NXProject is")
add_rect(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.5), AZUL_MEDIO)
add_textbox(slide, Inches(0.8), Inches(1.65), Inches(11.7), Inches(1.2),
            "A planning layer on top of Azure DevOps that turns the technical backlog "
            "into a manageable schedule — with dates, dependencies, Gantt, resource allocation "
            "and bidirectional traceability — without changing the technical team's workflow.",
            14, color=BRANCO, align=PP_ALIGN.CENTER)
pilares = [
    ("🗓️", "Automatic\nschedule",
     "Dates calculated from duration, work calendar, holidays and predecessors."),
    ("🔗", "Bidirectional\nintegration",
     "Imports and syncs with Azure DevOps: work items, sprints, owners and estimates."),
    ("📊", "Interactive\nGantt",
     "Bars, dependencies, milestones, critical path and baseline, with adjustable zoom."),
    ("🗂️", "Execution\nTaskBoard",
     "Sprint Tasks by person and state; the team runs the day and the board writes to DevOps."),
    ("👥", "Resource\nmanagement",
     "Allocation by person, overload detection, allocation map and cost per resource."),
    ("🤖", "Local or\ncloud AI",
     "Suggests task structures; the local model runs on your machine, sending no data out."),
]
for i, (emoji, titulo, corpo) in enumerate(pilares):
    l = Inches(0.2 + i * 2.17); t = Inches(3.25); w = Inches(2.05); h = Inches(3.8)
    add_rect(slide, l, t, w, h, BRANCO)
    add_textbox(slide, l, t+Inches(0.2), w, Inches(0.5), emoji, 26, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.1), t+Inches(0.75), w-Inches(0.2), Inches(0.75),
                titulo, 13, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.15), t+Inches(1.6), w-Inches(0.3), Inches(1.9),
                corpo, 10.5, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# ── Slide 4 — Philosophy: degrees of freedom ────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Planning philosophy: degrees of freedom")
add_rect(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.35), AZUL_MEDIO)
add_textbox(slide, Inches(0.8), Inches(1.62), Inches(11.7), Inches(1.1),
            "NXProject plans down to the Story level. Tasks can be pulled into the schedule, "
            "but the intent is for the team to detail and create them during execution, supported by the TaskBoard "
            "— which brings more agility to the project.",
            14, color=BRANCO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.6),
            "Inspired by the mathematical concept of degrees of freedom: define the boundaries of the system "
            "without constraining the movement inside them.",
            12, italic=True, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)
camadas = [
    ("Management sets the boundaries",
     "Dates, resources, dependencies and priority — at Epic, Feature and Story level.\nThat is what holds deadline, capacity and risk together."),
    ("The team moves inside them",
     "Tasks are born in execution, with the people doing the work: breakdown, sequence\nand fine tuning happen in the TaskBoard, with no replanning round trip."),
    ("Both sides see the same thing",
     "The Task created by the team rolls up to the Story in the schedule; managers see\nprogress and blockers without chasing status manually."),
]
for i, (titulo, corpo) in enumerate(camadas):
    l = Inches(0.5 + i * 4.2); t = Inches(3.95); w = Inches(3.9); h = Inches(2.6)
    add_rect(slide, l, t, w, h, BRANCO)
    add_rect(slide, l, t, w, Inches(0.06), LARANJA)
    add_textbox(slide, l+Inches(0.2), t+Inches(0.25), w-Inches(0.4), Inches(0.5),
                titulo, 13, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.2), t+Inches(1.0), w-Inches(0.4), Inches(1.5),
                corpo, 11, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# ── Slide 5 — Features ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Key features")
funcs = [
    ("Azure DevOps import",
     "• Project → Epic → Feature → Story hierarchy\n• Optional Task level (Load Task ToDo)\n• Estimates, dates, sprints and owners\n• Predecessors and blockers (Block tag)"),
    ("Smart schedule",
     "• Automatic predecessor cascade\n• Virtual predecessor by resource\n• Calendar with holidays and working hours\n• %Daily column: expected pace by today"),
    ("Gantt chart",
     "• Bars, milestones and dependency arrows\n• Critical path (CPM) and baseline\n• Zoom: Day, Week, Sprint, Month, Quarter\n• Drag to visually replan"),
    ("TaskBoard",
     "• Project & Story and Person & Task views\n• Drag to change state; batched write-back\n• Doing/Done and WIP limit per person\n• Alert for Stories out of sync"),
    ("Task Plan (Excel)",
     "• Breaks the Story down into Tasks\n• AI suggestions and Tech Lead review\n• Applies to the schedule and syncs back\n• Automatic backup on every save"),
    ("Resources and costs",
     "• Allocation by person and sprint\n• Allocation map by project and month\n• Cost per resource (hourly or monthly)\n• Overload alert (>100%)"),
    ("Bidirectional sync",
     "• Dates, hours, state, sprint, tags and predecessors\n• Creates new work items in DevOps\n• Concurrency control across users\n• Project admin group (Adm_NX)"),
    ("Health Check",
     "• Late activities (Finish < today and % < 100)\n• %Daily vs reported % deviation\n• Items without owner and circular dependencies"),
]
for i, (titulo, corpo) in enumerate(funcs):
    col = i % 4; row = i // 4
    l = Inches(0.35 + col * 3.25); t = Inches(1.55 + row * 2.85); w = Inches(3.05); h = Inches(2.6)
    add_rect(slide, l, t, w, h, BRANCO)
    add_rect(slide, l, t, w, Inches(0.45), AZUL_MEDIO)
    add_textbox(slide, l+Inches(0.15), t+Inches(0.05), w-Inches(0.2), Inches(0.38),
                titulo, 12, bold=True, color=BRANCO)
    add_textbox(slide, l+Inches(0.15), t+Inches(0.52), w-Inches(0.25), h-Inches(0.6),
                corpo, 10.5, color=CINZA_MEDIO)

# ── Slide 6 — Management benefits ─────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Strategic benefits for Management")
vantagens = [
    ("🎯  Real visibility into IT projects",
     "A consolidated schedule with dates, dependencies and completion percentage — without relying on manual team reports."),
    ("⚡  Zero impact on technical workflow",
     "The team continues using Azure DevOps as before. NXProject is a reading and planning layer — no process change required."),
    ("💰  Less rework and hidden cost",
     "Eliminates parallel spreadsheets, long status meetings and the cost of delays detected too late."),
    ("📈  Data-driven decisions",
     "Automatic Health Check, resource allocation and proactive delay alerts give management precise information for fast decisions."),
    ("🔒  Security and control",
     "Azure DevOps credentials protected by DPAPI (Windows per-user encryption). Project data stays in a local file — no mandatory cloud. Local AI runs the model on your own machine, sending no project data out."),
    ("📦  Simple deployment",
     "Desktop application with its own installer: no server, no database and no extra infrastructure. Ready in minutes."),
]
for i, (titulo, corpo) in enumerate(vantagens):
    col = i % 2; row = i // 2
    l = Inches(0.4 + col * 6.4); t = Inches(1.55 + row * 1.9); w = Inches(6.1); h = Inches(1.75)
    add_rect(slide, l, t, w, h, BRANCO)
    add_rect(slide, l, t, Inches(0.07), h, LARANJA)
    tb = slide.shapes.add_textbox(l+Inches(0.22), t+Inches(0.12), w-Inches(0.35), h-Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; run = p.add_run(); run.text = titulo
    run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = AZUL_ESCURO
    add_para(tf, corpo, 10.5, color=CINZA_MEDIO, space_before=5)

# ── Slide 7 — Workflow ────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "How it works — workflow")
etapas = [
    ("1", "Import",    "Connect to Azure DevOps\nand import the work item\nhierarchy"),
    ("2", "Plan",      "Adjust dates, duration,\nresources and dependencies\nin the schedule"),
    ("3", "Visualize", "Track dates, milestones\nand resource allocation\nin the Gantt"),
    ("4", "Execute",   "The team details the Tasks\nand runs the day to day\nin the TaskBoard"),
    ("5", "Sync",      "Send dates, hours, states\nand new work items back\nto DevOps"),
    ("6", "Monitor",   "Automatic Health Check\nalerts delays, blockers\nand overloads"),
]
step_w = Inches(1.9); arrow_w = Inches(0.3)
total_w = len(etapas) * step_w + (len(etapas) - 1) * arrow_w
start_l = (W - total_w) / 2
t_box = Inches(2.2); h_box = Inches(3.8)
for i, (num, titulo, corpo) in enumerate(etapas):
    l = start_l + i * (step_w + arrow_w)
    add_rect(slide, l, t_box, step_w, h_box, AZUL_MEDIO)
    add_rect(slide, l+Inches(0.6), t_box+Inches(0.18), Inches(0.7), Inches(0.7), LARANJA)
    add_textbox(slide, l+Inches(0.6), t_box+Inches(0.18), Inches(0.7), Inches(0.7),
                num, 20, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.1), t_box+Inches(1.1), step_w-Inches(0.2), Inches(0.6),
                titulo, 14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.1), t_box+Inches(1.85), step_w-Inches(0.2), Inches(1.8),
                corpo, 11, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
    if i < len(etapas) - 1:
        ax = l + step_w + Inches(0.1); ay = t_box + h_box/2 - Inches(0.15)
        add_textbox(slide, ax, ay, arrow_w-Inches(0.1), Inches(0.3),
                    "→", 22, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)

# ── Slide 8 — Who it is for ───────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
header(slide, "Who NXProject is for")
perfis = [
    ("🧑‍💼", "Project\nManager",
     "Schedule integrated with the backlog, delay alerts, dependency view and negotiated dates."),
    ("🔄", "Scrum Master\n/ RTE",
     "Sprint capacity, allocation conflicts, change impact and automatic cascade."),
    ("💻", "Tech Lead",
     "Feature and Story view with predecessors, hour estimates and traceability."),
    ("📊", "Management\n/ PMO",
     "Consolidated portfolio view, export to MS Project/Excel and executive Health Check."),
]
for i, (emoji, titulo, corpo) in enumerate(perfis):
    l = Inches(0.5 + i * 3.2); t = Inches(1.7); w = Inches(2.9); h = Inches(5.1)
    add_rect(slide, l, t, w, h, BRANCO)
    add_rect(slide, l, t, w, Inches(0.06), LARANJA)
    add_textbox(slide, l, t+Inches(0.2), w, Inches(0.7), emoji, 34, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.1), t+Inches(1.0), w-Inches(0.2), Inches(0.8),
                titulo, 14, bold=True, color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_textbox(slide, l+Inches(0.15), t+Inches(1.9), w-Inches(0.3), Inches(3.0),
                corpo, 11, color=CINZA_MEDIO, align=PP_ALIGN.CENTER)

# ── Slide 9 — Closing ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, AZUL_ESCURO)
add_rect(slide, 0, Inches(3.5), W, Inches(0.07), LARANJA)
add_textbox(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.0),
            "Turn your DevOps backlog into a manageable schedule.",
            26, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.1),
            "NXProject connects management and development — with transparency, "
            "automation and traceability — without changing the technical process.",
            16, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.7),
            "Request a demo", 22, bold=True, color=LARANJA, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.75), Inches(11.3), Inches(0.5),
            "comercial.nexus.xdata@gmail.com  •  nexusxdata.com.br",
            14, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
logo_h = Inches(1.25)
slide.shapes.add_picture(LOGO, (W - logo_h) / 2, Inches(5.55), height=logo_h)

# ── Salvar ────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "NXProject_Intelligent_DevOps_Planning_EN.pptx")
prs.save(out)
print(f"PPT saved to: {out}")
